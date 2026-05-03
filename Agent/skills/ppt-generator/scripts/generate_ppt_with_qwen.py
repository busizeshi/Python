import argparse
import json
import os
from pathlib import Path

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt

SYSTEM_PROMPT = """你是专业PPT策划顾问。请根据用户输入生成严格 JSON。
返回格式必须是 JSON 对象，字段如下：
{
  "deck_title": str,
  "subtitle": str,
  "audience": str,
  "duration_minutes": int,
  "tone": str,
  "language": str,
  "slides": [
    {
      "slide_no": int,
      "slide_title": str,
      "objective": str,
      "bullet_points": [str, str, str],
      "speaker_notes": str,
      "visual_suggestion": str
    }
  ]
}
要求：
1) 根据时长合理控制页数：10分钟6-8页，15分钟8-10页，30分钟12-16页。
2) 每页 bullet_points 3-5 条。
3) 内容可直接用于商务演示。
4) 不要输出 Markdown 或额外说明，只输出 JSON。
"""


def _load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def build_user_prompt(req: dict) -> str:
    return json.dumps(
        {
            "topic": req.get("topic"),
            "audience": req.get("audience", "通用"),
            "duration_minutes": int(req.get("duration_minutes", 15)),
            "tone": req.get("tone", "专业"),
            "language": req.get("language", "zh"),
        },
        ensure_ascii=False,
    )


def call_qwen_outline(req: dict) -> dict:
    load_dotenv()
    api_key = os.getenv("DASHSCOPE_API_KEY")
    base_url = os.getenv("QWEN_BASE_URL", "https://dashscope.aliyuncs.com/compatible-mode/v1")
    model = os.getenv("QWEN_MODEL", "qwen-plus")

    if not api_key:
        raise RuntimeError("缺少 DASHSCOPE_API_KEY。请先配置 .env。")

    client = OpenAI(api_key=api_key, base_url=base_url)
    completion = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": build_user_prompt(req)},
        ],
        temperature=0.4,
        response_format={"type": "json_object"},
    )

    content = completion.choices[0].message.content
    data = json.loads(content)
    if "slides" not in data or not isinstance(data["slides"], list):
        raise ValueError("模型输出缺少 slides 数组")
    return data


def render_ppt(outline: dict, pptx_path: Path) -> None:
    prs = Presentation()

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = outline.get("deck_title", "PPT")
    subtitle = cover.placeholders[1] if len(cover.placeholders) > 1 else None
    if subtitle:
        subtitle.text = outline.get("subtitle", "")

    for slide in outline.get("slides", []):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide.get("slide_title", "")

        body = s.shapes.placeholders[1].text_frame
        body.clear()
        objective = slide.get("objective", "")
        if objective:
            p = body.paragraphs[0]
            p.text = f"目标：{objective}"
            p.font.bold = True
            p.font.size = Pt(18)

        for bp in slide.get("bullet_points", []):
            p = body.add_paragraph()
            p.text = str(bp)
            p.level = 1
            p.font.size = Pt(16)

        notes = s.notes_slide.notes_text_frame
        notes.text = (
            f"讲稿备注：{slide.get('speaker_notes', '')}\n"
            f"视觉建议：{slide.get('visual_suggestion', '')}"
        )

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", default="D:\\dev\\Python\\Agent\\skills\\ppt-generator\\examples\\request_sample.json")
    parser.add_argument("--outline-output", default="D:\\dev\\Python\\Agent\\skills\\ppt-generator\\examples\\outline_llm_output.json")
    parser.add_argument("--pptx-output", default="D:\\dev\\Python\\Agent\\skills\\ppt-generator\\examples\\slides_from_qwen.pptx")
    args = parser.parse_args()

    req = _load_json(Path(args.input))
    outline = call_qwen_outline(req)

    outline_path = Path(args.outline_output)
    outline_path.write_text(json.dumps(outline, ensure_ascii=False, indent=2), encoding="utf-8")

    render_ppt(outline, Path(args.pptx_output))

    print(f"Outline generated -> {outline_path}")
    print(f"PPT generated -> {args.pptx_output}")


if __name__ == "__main__":
    main()
